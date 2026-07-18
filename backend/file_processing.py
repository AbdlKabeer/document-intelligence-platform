from typing import List, Tuple
import pdfplumber, pytesseract, docx
import pandas as pd
from PIL import Image
import io
#from utils import process_table
from operations.table_operations import process_table
from schemas import Table
from fastapi import HTTPException
from pptx import Presentation
from config import logger

# File extraction functions
def extract_text_from_pdf(file_content: bytes) -> Tuple[str, List[Table]]:
    """Extract text and tables from PDF using pdfplumber."""
    pdf_file = io.BytesIO(file_content)
    text = ""
    tables_list = []
    table_counter = 0
    
    try:
        with pdfplumber.open(pdf_file) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text += page_text + "\n"
                
                tables = page.extract_tables()
                if tables:
                    for table_data in tables:
                        table_counter += 1
                        table = process_table(table_data, table_counter, page_text)
                        tables_list.append(table)
                        text += f"\n[Reference to {table.name}]\n"
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        raise HTTPException(status_code=500, detail="Error extracting text from PDF.")
    
    return text, tables_list

def extract_text_from_xlsx(file_content: bytes) -> Tuple[str, List[Table]]:
    """Extract text and tables from an Excel file."""
    excel_file = io.BytesIO(file_content)
    text = ""
    tables_list = []
    
    try:
        sheets = pd.read_excel(excel_file, sheet_name=None)
        for sheet_name, df in sheets.items():
            csv_content = df.to_csv(index=False)
            table = Table(
                name=f"sheet_{sheet_name}",
                content=csv_content,
                description=f"Data from sheet '{sheet_name}' containing columns: {', '.join(df.columns)}"
            )
            tables_list.append(table)
            text += f"\nSheet: {sheet_name} [Reference to {table.name}]\n"
            text += f"Contains columns: {', '.join(df.columns)}\n"
    except Exception as e:
        logger.error(f"Error processing Excel file: {e}")
        raise HTTPException(status_code=500, detail="Error processing Excel file.")
    
    return text, tables_list

def extract_text_from_docx(file_content: bytes) -> Tuple[str, List[Table]]:
    """Extract text and tables from a .docx file."""
    doc_file = io.BytesIO(file_content)
    document = docx.Document(doc_file)
    text = "\n".join([para.text for para in document.paragraphs])
    tables_list = []
    table_counter = 0
    
    for table in document.tables:
        table_counter += 1
        table_data = []
        for row in table.rows:
            table_data.append([cell.text for cell in row.cells])
        
        table = process_table(table_data, table_counter, text)
        tables_list.append(table)
        text += f"\n[Reference to {table.name}]\n"
    
    return text, tables_list

def extract_text_from_txt(file_content: bytes) -> Tuple[str, List[Table]]:
    """Extract text from a .txt file."""
    return file_content.decode("utf-8"), []

def extract_text_from_image(image_content) -> Tuple[str, List[Table]]:
    """Extract text from scanned documents using Tesseract OCR."""
    image = Image.open(image_content)
    return pytesseract.image_to_string(image), []

def extract_text_from_pptx(file_content: bytes) -> Tuple[str, List[Table]]:
    """Extract text and tables from PowerPoint presentations."""
    pptx_file = io.BytesIO(file_content)
    text = ""
    tables_list = []
    table_counter = 0
    
    try:
        presentation = Presentation(pptx_file)
        
        for slide_number, slide in enumerate(presentation.slides, 1):
            text += f"\n--- Slide {slide_number} ---\n"
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                
                # Extract tables
                if shape.has_table:
                    table_counter += 1
                    table_data = []
                    
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            # Get text from each cell, handling potential formatting
                            cell_text = ""
                            for paragraph in cell.text_frame.paragraphs:
                                cell_text += paragraph.text + " "
                            row_data.append(cell_text.strip())
                        table_data.append(row_data)
                    
                    # Process the table using the existing utility
                    table = process_table(table_data, table_counter, text)
                    tables_list.append(table)
                    text += f"\n[Reference to {table.name}]\n"
                
                # Extract text from text boxes and other shapes with text frames
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + " "
                        text += "\n"
            
            # Extract notes if present
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    text += "\nSlide Notes:\n"
                    text += notes_slide.notes_text_frame.text + "\n"
    
    except Exception as e:
        logger.error(f"Error extracting text from PowerPoint: {e}")
        raise HTTPException(status_code=500, detail="Error extracting text from PowerPoint presentation.")
    
    return text, tables_list
